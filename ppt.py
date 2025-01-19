from pptx import Presentation

# Create a presentation object
presentation = Presentation()

# Function to add a slide with a title and bullet points
def add_slide_with_bullets(title, bullets):
    slide_layout = presentation.slide_layouts[1]  # Title and Content layout
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    for bullet in bullets:
        content.text += f"- {bullet}\n"

# Title Slide
title_slide_layout = presentation.slide_layouts[0]
title_slide = presentation.slides.add_slide(title_slide_layout)
title_slide.shapes.title.text = "Headache Decision Support System (*HeadacheDSS*)"
title_slide.placeholders[1].text = (
    '"Diagnosing and Monitoring Chronic Primary Headaches"\n\nPresented By: [Your Name]\nDate: [Presentation Date]'
)

# Slides Content
slides_content = [
    ("Introduction", [
        "What is HeadacheDSS?",
        "A Decision Support System for diagnosing chronic primary headaches.",
        "Combines ICHD standards and machine learning.",
        "Why HeadacheDSS?",
        "Addresses limitations of current diagnostic methods.",
        "Key Feature: Utilizes knowledge graphs for better predictions."
    ]),
    ("Objectives", [
        "Primary Goal: Build a knowledge-based system (KBS) for headache diagnosis.",
        "Develop a knowledge base using ICHD guidelines.",
        "Apply semantic enrichment for enhanced feature extraction.",
        "Evaluate oversampling techniques for imbalanced datasets.",
        "Train and assess machine learning models."
    ]),
    ("Problem Description", [
        "Challenges in Headache Diagnosis:",
        "Similar symptoms lead to misdiagnosis.",
        "Imbalanced datasets for uncommon headache types.",
        "Proposed Solution:",
        "Enhance clinical datasets using knowledge graphs.",
        "Use machine learning for accurate classification."
    ]),
    ("Methodology Overview", [
        "Step-by-step process:",
        "Dataset Preparation",
        "Knowledge Base Construction",
        "Semantic Enrichment",
        "Oversampling Techniques",
        "Feature Extraction",
        "Machine Learning Model Training"
    ]),
    ("Results", [
        "Key Findings:",
        "Enriched features improved classification accuracy.",
        "Domain-informed oversampling outperformed SMOTE and ADASYN.",
        "Model Evaluation:",
        "Decision Tree with enriched features achieved the best results."
    ]),
    ("Conclusion", [
        "Summary:",
        "HeadacheDSS combines structured knowledge and ML for better diagnosis.",
        "Semantic enrichment and domain-informed sampling enhance prediction accuracy.",
        "Next Steps:",
        "Expand the dataset.",
        "Experiment with advanced ML models (e.g., Random Forest, SVM).",
        "Further refine the knowledge base."
    ]),
]

# Add slides to the presentation
for title, bullets in slides_content:
    add_slide_with_bullets(title, bullets)

# Add Thank You Slide
thank_you_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
thank_you_slide.shapes.title.text = "Thank You"
thank_you_slide.placeholders[1].text = "Questions?\n\n[Your Contact Information or Institution Name]"

# Save the presentation
presentation.save("HeadacheDSS_Presentation.pptx")
